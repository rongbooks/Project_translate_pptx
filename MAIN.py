import os
import tkinter as tk
from tkinter import filedialog, messagebox, ttk
import threading
from pptx import Presentation
import requests
import json
import hashlib
import time
import random


class BaiduTranslator:
    def __init__(self, appid, secret_key):
        self.appid = appid
        self.secret_key = secret_key
        self.api_url = 'https://fanyi-api.baidu.com/api/trans/vip/fieldtranslate'

    def translate(self, text, from_lang='en', to_lang='zh'):
        if not text.strip():
            return text

        salt = random.randint(32768, 65536)
        sign = self.appid + text + str(salt) + self.secret_key
        sign = hashlib.md5(sign.encode()).hexdigest()

        try:
            response = requests.get(self.api_url, params={
                'q': text,
                'from': from_lang,
                'to': to_lang,
                'appid': self.appid,
                'salt': salt,
                'sign': sign
            })
            result = response.json()
            if 'trans_result' in result:
                return result['trans_result'][0]['dst']
            else:
                print(f"翻译失败: {result.get('error_msg', '未知错误')}")
                return text
        except Exception as e:
            print(f"翻译请求异常: {e}")
            return text


class PPTTranslatorApp:
    def __init__(self, root):
        self.root = root
        self.root.title("PPT翻译工具")
        self.root.geometry("600x400")
        self.root.resizable(True, True)

        # 设置中文字体
        self.style = ttk.Style()
        self.style.configure("TLabel", font=("SimHei", 10))
        self.style.configure("TButton", font=("SimHei", 10))
        self.style.configure("TProgressbar", thickness=20)

        # 创建界面
        self.create_widgets()

        # 百度翻译API信息
        self.appid = ""
        self.secret_key = ""

        # 翻译器
        self.translator = None

    def create_widgets(self):
        # 主框架
        main_frame = ttk.Frame(self.root, padding="20")
        main_frame.pack(fill=tk.BOTH, expand=True)

        # API设置部分
        api_frame = ttk.LabelFrame(main_frame, text="百度翻译API设置", padding="10")
        api_frame.pack(fill=tk.X, pady=10)

        ttk.Label(api_frame, text="APP ID:").grid(row=0, column=0, sticky=tk.W, pady=5)
        self.appid_var = tk.StringVar()
        ttk.Entry(api_frame, textvariable=self.appid_var, width=40).grid(row=0, column=1, sticky=tk.W, pady=5)

        ttk.Label(api_frame, text="Secret Key:").grid(row=1, column=0, sticky=tk.W, pady=5)
        self.secret_key_var = tk.StringVar()
        ttk.Entry(api_frame, textvariable=self.secret_key_var, show="*", width=40).grid(row=1, column=1, sticky=tk.W,
                                                                                        pady=5)

        # 文件选择部分
        file_frame = ttk.LabelFrame(main_frame, text="文件选择", padding="10")
        file_frame.pack(fill=tk.X, pady=10)

        self.file_path_var = tk.StringVar()
        ttk.Entry(file_frame, textvariable=self.file_path_var, state="readonly", width=40).grid(row=0, column=0,
                                                                                                sticky=tk.W, pady=5)
        ttk.Button(file_frame, text="浏览", command=self.browse_file).grid(row=0, column=1, padx=5, pady=5)

        # 翻译按钮
        button_frame = ttk.Frame(main_frame)
        button_frame.pack(pady=20)

        self.translate_button = ttk.Button(button_frame, text="开始翻译", command=self.start_translation)
        self.translate_button.pack(padx=10, pady=10)

        # 进度条
        self.progress_var = tk.DoubleVar()
        self.progress_bar = ttk.Progressbar(main_frame, variable=self.progress_var, maximum=100)
        self.progress_bar.pack(fill=tk.X, pady=10)

        # 状态标签
        self.status_var = tk.StringVar()
        self.status_var.set("就绪")
        ttk.Label(main_frame, textvariable=self.status_var).pack(anchor=tk.W, pady=5)

    def browse_file(self):
        file_path = filedialog.askopenfilename(
            title="选择PPT文件",
            filetypes=[("PowerPoint文件", "*.pptx")]
        )
        if file_path:
            self.file_path_var.set(file_path)

    def start_translation(self):
        # 检查API信息
        self.appid = self.appid_var.get().strip()
        self.secret_key = self.secret_key_var.get().strip()

        if not self.appid or not self.secret_key:
            messagebox.showerror("错误", "请输入百度翻译API的APP ID和Secret Key")
            return

        # 检查文件
        file_path = self.file_path_var.get()
        if not file_path or not os.path.exists(file_path):
            messagebox.showerror("错误", "请选择有效的PPT文件")
            return

        # 禁用按钮
        self.translate_button.config(state="disabled")
        self.status_var.set("正在翻译...")

        # 创建翻译器
        self.translator = BaiduTranslator(self.appid, self.secret_key)

        # 在新线程中执行翻译
        threading.Thread(target=self.perform_translation, daemon=True).start()

    def perform_translation(self):
        try:
            file_path = self.file_path_var.get()
            prs = Presentation(file_path)

            total_slides = len(prs.slides)
            translated_slides = 0

            # 创建文本记录文件
            base_name = os.path.splitext(file_path)[0]
            txt_file = base_name + "_翻译记录.txt"
            with open(txt_file, 'w', encoding='utf-8') as f:
                f.write(f"PPT文件: {os.path.basename(file_path)}\n")
                f.write(f"翻译时间: {time.strftime('%Y-%m-%d %H:%M:%S')}\n")
                f.write("=" * 50 + "\n\n")

            # 遍历所有幻灯片
            for slide in prs.slides:
                slide_number = prs.slides.index(slide) + 1

                # 遍历所有形状
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue

                    # 遍历文本框中的所有段落
                    for paragraph in shape.text_frame.paragraphs:
                        original_text = ""
                        for run in paragraph.runs:
                            original_text += run.text

                        # 只翻译英文文本
                        if self.is_english_text(original_text):
                            # 翻译文本
                            translated_text = self.translator.translate(original_text)

                            # 保存到文本文件
                            with open(txt_file, 'a', encoding='utf-8') as f:
                                f.write(f"幻灯片 #{slide_number}\n")
                                f.write(f"{original_text}\n")
                                f.write(f"{translated_text}\n\n")

                            # 保留格式替换文本
                            if paragraph.runs:
                                # 保留第一个运行的格式，替换其文本
                                paragraph.runs[0].text = translated_text
                                # 移除其余的运行
                                for i in range(len(paragraph.runs) - 1, 0, -1):
                                    p = paragraph._element
                                    p.remove(paragraph.runs[i]._r)

                # 更新进度
                translated_slides += 1
                progress = (translated_slides / total_slides) * 100
                self.root.after(100, lambda p=progress: self.progress_var.set(p))

            # 保存翻译后的PPT
            output_file = base_name + "_翻译版.pptx"
            prs.save(output_file)

            # 更新状态
            self.root.after(0, lambda: self.status_var.set(f"翻译完成！文件已保存为: {output_file}"))
            self.root.after(0, lambda: messagebox.showinfo("成功",
                                                           f"PPT翻译完成！\n"
                                                           f"PPT文件已保存为: {output_file}\n"
                                                           f"翻译记录已保存为: {txt_file}"
                                                           ))
        except Exception as e:
            self.root.after(0, lambda: self.status_var.set(f"翻译失败: {str(e)}"))
            self.root.after(0, lambda: messagebox.showerror("错误", f"翻译过程中发生错误: {str(e)}"))
        finally:
            # 重新启用按钮
            self.root.after(0, lambda: self.translate_button.config(state="normal"))

    def is_english_text(self, text):
        """简单判断文本是否主要为英文"""
        if not text:
            return False

        english_chars = sum(1 for c in text if c.isascii())
        total_chars = len(text.strip())

        # 如果文本中超过50%的字符是ASCII字符，则认为是英文
        return (english_chars / total_chars) > 0.5 if total_chars > 0 else False


if __name__ == "__main__":
    root = tk.Tk()
    app = PPTTranslatorApp(root)
    root.mainloop()